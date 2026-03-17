import fitz 
import sys
import pytesseract
from pptx import Presentation
from docx import Document
from PIL import Image

def read_pdf(path):

    full_text = ""
    try:
        with fitz.open(path) as pdf:
            for page in pdf:
                full_text += page.get_text().strip()
        return full_text
    
    except Exception as e:
        print("Error opening or reading pdf: " + str(e))
        return None
    
def read_pptx(path):
    try:
        pptx = Presentation(path)

        text = []

        for slide in pptx.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text.append(shape.text.strip())

        # Combine all strings separate by new lines
        return "\n".join(text)

    except Exception as e:
        print("Error opening or reading pptx: " + str(e))
        return None
    
def read_docx(path):
    try: 
        doc = Document(path)
        text = []

        for para in doc.paragraphs:
            text.append(para.text.strip())

        return "\n".join(text)

    except Exception as e:
        print("Error opening or reading docx: " + str(e))
        return None

def read_ocr(path):
    img = Image.open(path)
    text = pytesseract.image_to_string(img)

    return text


if __name__ == "__main__":

    # 1st arg is .py file, 2nd is file to be read
    path = sys.argv[1]

    if(path.endswith(".pdf")):
        print(read_pdf(path))
    elif(path.endswith(".pptx")):
        print(read_pptx(path))
    elif(path.endswith(".docx")):
        print(read_docx(path))
    elif(path.endswith(".png") or path.endswith(".jpg") or path.endswith(".jpeg")):
        print(read_ocr(path))
    else:
        print("Unknown file type")
