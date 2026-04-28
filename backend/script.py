import fitz 
import sys
import pytesseract

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io

from docx import Document
from PIL import Image

import re

def read_pdf(path, start_page=None, end_page=None):

    full_text = ""
    try:
        with fitz.open(path) as pdf:
            total_pages = len(pdf)
            start_idx = 0
            end_idx = total_pages - 1

            if start_page is not None or end_page is not None:
                if start_page is None or end_page is None:
                    raise ValueError("Both start and end pages are required for a custom range.")
                if start_page < 1 or end_page < 1:
                    raise ValueError("Page numbers must be 1 or greater.")
                if start_page > end_page:
                    raise ValueError("Start page cannot be greater than end page.")
                if end_page > total_pages:
                    raise ValueError(f"This PDF has {total_pages} pages. End page exceeds that.")
                start_idx = start_page - 1
                end_idx = end_page - 1
            
            # Blocks preserve original PDF paragraph format
            for page_num in range(start_idx, end_idx + 1):
                page = pdf[page_num]
                blocks = page.get_text("blocks")
                blocks.sort(key=lambda block: (block[1], block[0]))   # Sort blocks from top to bottom (y to x) with lambda function, default is left to right (x to y)
                
                for block in blocks:
                    block_text = block[4].strip()   # Gets block text
                    if not block_text:
                        continue

                    full_text += block_text + "\n\n"
                    
                
        return full_text.strip()
    
    except ValueError:
        raise
    except Exception as e:
        print("Error opening or reading pdf: " + str(e))
        return None
    
def read_pptx(path):
    try:
        pptx = Presentation(path)

        text = []

        for slide in pptx.slides:
            for shape in slide.shapes:

                if shape.has_text_frame:
                    text.append(shape.text.strip())

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:

                    # Raw bytes
                    img_bytes = shape.image.blob

                    img_text = read_ocr_bytes(img_bytes)

                    text.append(img_text.strip())

        # Combine all strings separate by new lines
        return "\n".join(text)

    except Exception as e:
        print("Error opening or reading pptx: " + str(e))
        return None
    
def read_docx(path):
    try: 
        doc = Document(path)
        text = []

        for para in doc.paragraphs:
            text.append(para.text.strip())

        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                img_bytes = rel.target_part.blob
                img_text = read_ocr_bytes(img_bytes)
                if img_text:
                    text.append(img_text.strip())

        return "\n".join(text)

    except Exception as e:
        print("Error opening or reading docx: " + str(e))
        return None

def read_ocr_path(path):
    img = Image.open(path)
    text = pytesseract.image_to_string(img)

    return text

def read_ocr_bytes(bytes):
    img = Image.open(io.BytesIO(bytes))

    # Assumes block of text (--psm 6)
    text = pytesseract.image_to_string(img, config="--oem 3 --psm 6")

    # Regex: text has to be letters/numbers, with 2 or more characters
    if not re.search(r"[A-Za-z0-9]{2,}", text):
        return ""

    return text


if __name__ == "__main__":

    # 1st arg is .py file, 2nd is file to be read
    path = sys.argv[1]

    if(path.endswith(".pdf")):
        print(read_pdf(path))
    elif(path.endswith(".pptx")):
        print(read_pptx(path))
    elif(path.endswith(".docx")):
        print(read_docx(path))
    elif(path.endswith(".png") or path.endswith(".jpg") or path.endswith(".jpeg")):
        print(read_ocr_path(path))
    else:
        print("Unknown file type")
